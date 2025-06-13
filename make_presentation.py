from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    body = slide.placeholders[1].text_frame
    for point in bullet_points:
        p = body.add_paragraph() if body.text else body.paragraphs[0]
        p.text = point

add_slide("Colony Detection SAM 3.0", ["基于SAM的菌落检测和分析", "自动解析文件名与结果组织", "支持96孔板自动对齐和分析"])
add_slide("代码结构", ["main.py: 命令行入口", "pipeline.py: 核心流程", "core/: 检测与分析模块"])
add_slide("运行流程", ["加载配置并初始化模型", "检测菌落（auto/grid/hybrid）", "分析特征并保存报告"])
add_slide("调试与可视化", ["--debug 可生成中间结果", "integrate_report.py 汇总生成报告"])
add_slide("近期修复", ["新增 json import", "清理重复的 self.args 赋值"])

prs.save('Colony_Detection_Overview.pptx')
print('ppt saved')
